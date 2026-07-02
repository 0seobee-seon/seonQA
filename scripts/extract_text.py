"""
Phase 2: 텍스트 추출 스크립트
knowledge_base 폴더의 모든 문서에서 텍스트를 추출하여 JSON으로 저장
"""

import os, json, re, zlib, struct, zipfile
import olefile
import pdfplumber
import win32com.client
from pptx import Presentation
from pathlib import Path

BASE_DIR   = Path(r"C:\Users\LG\Desktop\AX 코드 학습\챗봇데이터\knowledge_base")
OUTPUT_DIR = Path(r"C:\Users\LG\Desktop\AX 코드 학습\챗봇데이터\extracted")
OUTPUT_DIR.mkdir(exist_ok=True)


# ── 추출 함수 ────────────────────────────────────────────────────

def extract_hwp(filepath: Path) -> str:
    """HWP (OLE 바이너리) 파일에서 텍스트 추출"""
    try:
        with olefile.OleFileIO(str(filepath)) as ole:
            sections = [s for s in ole.listdir() if s[0] == "BodyText"]
            if not sections:
                return ""
            lines = []
            for sec in sections:
                data = ole.openstream("/".join(sec)).read()
                try:
                    data = zlib.decompress(data, -15)
                except Exception:
                    pass
                i = 0
                while i < len(data) - 4:
                    tag = struct.unpack_from("<I", data, i)[0]
                    rec_type = tag & 0x3FF
                    rec_len  = (tag >> 20) & 0xFFF
                    if rec_len == 0xFFF:
                        rec_len = struct.unpack_from("<I", data, i + 4)[0]
                        i += 4
                    i += 4
                    body = data[i : i + rec_len]
                    if rec_type == 67:  # HWPTAG_PARA_TEXT
                        try:
                            t = body.decode("utf-16-le").replace("\x00", "").strip()
                            if t:
                                lines.append(t)
                        except Exception:
                            pass
                    i += rec_len
            return _clean(lines)
    except Exception as e:
        return f"[추출 실패: {e}]"


def extract_hwpx(filepath: Path) -> str:
    """HWPX (ZIP+XML) 파일에서 텍스트 추출"""
    try:
        with zipfile.ZipFile(str(filepath)) as z:
            names = z.namelist()
            # section XML 파일 전부 처리
            sections = sorted([n for n in names if re.match(r"Contents/section\d+\.xml", n)])
            lines = []
            for sec in sections:
                xml = z.read(sec).decode("utf-8", errors="ignore")
                texts = re.findall(r"<hp:t[^>]*>([^<]+)</hp:t>", xml)
                lines.extend(texts)
            return _clean(lines)
    except Exception as e:
        return f"[추출 실패: {e}]"


def extract_pdf(filepath: Path) -> str:
    """PDF 파일에서 텍스트 추출 (원본 PDF 전용)"""
    try:
        lines = []
        with pdfplumber.open(str(filepath)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    lines.append(t)
        return "\n".join(lines).strip()
    except Exception as e:
        return f"[추출 실패: {e}]"


def extract_ppt(filepath: Path) -> str:
    """구형 PPT 파일에서 텍스트 추출 (PowerPoint COM)"""
    ppt_app = None
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = 1
        prs = ppt_app.Presentations.Open(
            str(filepath), ReadOnly=True, Untitled=False, WithWindow=False
        )
        lines = []
        for slide in prs.Slides:
            slide_texts = []
            for shape in slide.Shapes:
                try:
                    if shape.HasTextFrame:
                        t = shape.TextFrame.TextRange.Text.strip()
                        if t:
                            slide_texts.append(t)
                except Exception:
                    pass
            if slide_texts:
                lines.append(f"[슬라이드 {slide.SlideIndex}]")
                lines.extend(slide_texts)
        prs.Close()
        return "\n".join(lines).strip()
    except Exception as e:
        return f"[추출 실패: {e}]"
    finally:
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass


def extract_pptx(filepath: Path) -> str:
    """PPTX 파일에서 텍스트 추출"""
    try:
        prs = Presentation(str(filepath))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                lines.append(f"[슬라이드 {i}]")
                lines.extend(slide_texts)
        return "\n".join(lines).strip()
    except Exception as e:
        return f"[추출 실패: {e}]"


def extract_md(filepath: Path) -> str:
    try:
        return filepath.read_text(encoding="utf-8").strip()
    except Exception as e:
        return f"[추출 실패: {e}]"


def extract_csv(filepath: Path) -> str:
    """CSV → 자연어 문장 변환 (연락처 특화)"""
    try:
        import csv
        rows = []
        with open(filepath, encoding="utf-8-sig", newline="") as f:
            reader = csv.DictReader(f)
            for row in reader:
                if not any(row.values()):
                    continue
                parts = [f"{k}: {v}" for k, v in row.items() if v and v.strip()]
                if parts:
                    rows.append(", ".join(parts))
        return "\n".join(rows).strip()
    except Exception as e:
        return f"[추출 실패: {e}]"


def _clean(lines: list) -> str:
    """추출된 텍스트 정제 (깨진 문자 제거)"""
    text = "\n".join(lines)
    # 한글/영문/숫자/기본 특수문자 외 제거
    text = re.sub(r"[^가-힣ᄀ-ᇿ㄰-㆏"
                  r"a-zA-Z0-9 \-_.,!?()\[\]{}/\\:;@#%&=+*\n\r\t\"'<>~^]+",
                  " ", text)
    text = re.sub(r" {2,}", " ", text)
    return text.strip()


# ── 카테고리 추론 ─────────────────────────────────────────────────

CATEGORY_MAP = {
    "규정정책":           "규정정책",
    "업무메뉴얼":         "업무매뉴얼",
    "사내 연락처":        "조직연락처",
    "사내 양식":          "사내양식",
    "bizmeka groupware":  "업무매뉴얼",
}

def get_category(filepath: Path) -> str:
    for part in filepath.parts:
        if part in CATEGORY_MAP:
            return CATEGORY_MAP[part]
    return "기타"


# ── 형식별 추출 함수 매핑 ──────────────────────────────────────────

EXTRACT_MAP = {
    ".hwp":  extract_hwp,
    ".hwpx": extract_hwpx,
    ".ppt":  extract_ppt,
    ".pptx": extract_pptx,
    ".md":   extract_md,
    ".csv":  extract_csv,
}

# PDF는 원본 HWP 없는 폴더만 처리
PDF_ONLY_FOLDERS = {"bizmeka groupware"}


# ── 메인 처리 ─────────────────────────────────────────────────────

results = []
skipped = []

all_files = [f for f in BASE_DIR.rglob("*") if f.is_file()]
print(f"총 파일 수: {len(all_files)}\n")

for filepath in all_files:
    ext = filepath.suffix.lower()
    folder_top = filepath.relative_to(BASE_DIR).parts[0] if filepath.relative_to(BASE_DIR).parts else ""

    # HWP 변환 PDF는 건너뜀
    if ext == ".pdf":
        if folder_top not in PDF_ONLY_FOLDERS:
            skipped.append(filepath.name)
            continue
        extractor = extract_pdf
    elif ext in EXTRACT_MAP:
        extractor = EXTRACT_MAP[ext]
    else:
        continue

    print(f"처리 중: {filepath.name}")
    content = extractor(filepath)

    if not content or content.startswith("[추출 실패"):
        print(f"  [!] 실패: {filepath.name} | {content[:80] if content else '내용 없음'}")
        continue

    doc = {
        "title":       filepath.stem,
        "category":    get_category(filepath),
        "source_file": str(filepath.relative_to(BASE_DIR)),
        "content":     content,
    }
    results.append(doc)

    safe_name = re.sub(r'[\\/:*?"<>|]', "_", filepath.stem)
    out_path = OUTPUT_DIR / f"{safe_name}.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(doc, f, ensure_ascii=False, indent=2)

# 전체 통합 JSON
all_path = OUTPUT_DIR / "_all_documents.json"
with open(all_path, "w", encoding="utf-8") as f:
    json.dump(results, f, ensure_ascii=False, indent=2)

print(f"\n완료: {len(results)}개 문서 추출 -> {OUTPUT_DIR}")
print(f"건너뜀 (HWP 변환 PDF): {len(skipped)}개")
